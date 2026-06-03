"""Extract plain text from uploaded documents (PDF / images / PPT / office)."""

from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_DOCUMENT_BYTES = 20 * 1024 * 1024

TEXT_EXTENSIONS = frozenset({".txt", ".md", ".markdown", ".csv", ".json"})
PDF_EXTENSIONS = frozenset({".pdf"})
IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"})
PPTX_EXTENSIONS = frozenset({".pptx"})
PPT_EXTENSIONS = frozenset({".ppt"})
DOCX_EXTENSIONS = frozenset({".docx"})

SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | PDF_EXTENSIONS
    | IMAGE_EXTENSIONS
    | PPTX_EXTENSIONS
    | PPT_EXTENSIONS
    | DOCX_EXTENSIONS
)


class DocumentExtractError(Exception):
    def __init__(self, message: str, code: str = "extract_failed") -> None:
        super().__init__(message)
        self.code = code


@dataclass(frozen=True)
class DocumentExtractResult:
    text: str
    filename: str
    method: str
    char_count: int


_ocr_engine = None


def _get_ocr_engine():
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR

        _ocr_engine = RapidOCR()
    return _ocr_engine


def _ocr_image_bytes(data: bytes) -> str:
    from PIL import Image

    engine = _get_ocr_engine()
    image = Image.open(io.BytesIO(data))
    if image.mode not in ("RGB", "L"):
        image = image.convert("RGB")
    result, _ = engine(image)
    if not result:
        return ""
    lines: list[str] = []
    for item in result:
        if len(item) >= 2 and item[1]:
            lines.append(str(item[1]).strip())
    return "\n".join(lines)


def _normalize_text(text: str) -> str:
    cleaned = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _extract_plain_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentExtractError("无法解码文本文件编码", "decode_error")


def _extract_pdf(data: bytes) -> tuple[str, str]:
    import fitz
    import numpy as np

    doc = fitz.open(stream=data, filetype="pdf")
    parts: list[str] = []
    ocr_pages = 0
    try:
        for page in doc:
            text = (page.get_text("text") or "").strip()
            if len(text) >= 40:
                parts.append(text)
                continue

            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height,
                pix.width,
                pix.n,
            )
            if pix.n == 4:
                img = img[:, :, :3]
            engine = _get_ocr_engine()
            result, _ = engine(img)
            if result:
                ocr_pages += 1
                lines = [str(item[1]).strip() for item in result if len(item) >= 2]
                page_text = "\n".join(x for x in lines if x)
                if page_text:
                    parts.append(page_text)
            elif text:
                parts.append(text)
    finally:
        doc.close()

    merged = _normalize_text("\n\n".join(parts))
    if not merged:
        raise DocumentExtractError("PDF 中未识别到可读文字", "empty_document")
    method = "pdf_ocr" if ocr_pages else "pdf_text"
    if ocr_pages and len(merged) > 40:
        method = "pdf_mixed" if len(parts) > ocr_pages else "pdf_ocr"
    return merged, method


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for slide in prs.slides:
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                slide_lines.append(text)
        if slide_lines:
            chunks.append("\n".join(slide_lines))
    merged = _normalize_text("\n\n".join(chunks))
    if not merged:
        raise DocumentExtractError("PPT 中未找到可提取文字", "empty_document")
    return merged


def _extract_ppt_legacy(data: bytes) -> str:
    """Legacy .ppt via optional MarkItDown (requires extra install)."""
    try:
        from markitdown import MarkItDown
    except ImportError as e:
        raise DocumentExtractError(
            "旧版 .ppt 暂不支持，请另存为 .pptx 后上传",
            "unsupported_format",
        ) from e

    converter = MarkItDown()
    result = converter.convert_stream(io.BytesIO(data), file_extension=".ppt")
    text = _normalize_text(getattr(result, "text_content", "") or "")
    if not text:
        raise DocumentExtractError(
            "无法从 .ppt 提取文字，请另存为 .pptx",
            "empty_document",
        )
    return text


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    merged = _normalize_text("\n".join(parts))
    if not merged:
        raise DocumentExtractError("Word 文档中未找到可提取文字", "empty_document")
    return merged


def extract_document(data: bytes, filename: str) -> DocumentExtractResult:
    if len(data) > MAX_DOCUMENT_BYTES:
        raise DocumentExtractError(
            f"文件过大（上限 {MAX_DOCUMENT_BYTES // (1024 * 1024)} MB）",
            "file_too_large",
        )

    ext = Path(filename or "").suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        allowed = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise DocumentExtractError(
            f"不支持的文件类型 {ext or '(无扩展名)'}，支持：{allowed}",
            "unsupported_format",
        )

    method = "plain_text"
    if ext in TEXT_EXTENSIONS:
        text = _extract_plain_text(data)
    elif ext in PDF_EXTENSIONS:
        text, method = _extract_pdf(data)
    elif ext in IMAGE_EXTENSIONS:
        text = _ocr_image_bytes(data)
        method = "image_ocr"
    elif ext in PPTX_EXTENSIONS:
        text = _extract_pptx(data)
        method = "pptx"
    elif ext in PPT_EXTENSIONS:
        text = _extract_ppt_legacy(data)
        method = "ppt_legacy"
    elif ext in DOCX_EXTENSIONS:
        text = _extract_docx(data)
        method = "docx"
    else:
        raise DocumentExtractError("不支持的文件类型", "unsupported_format")

    text = _normalize_text(text)
    if len(text) < 1:
        raise DocumentExtractError("未能从文件中提取到文字", "empty_document")

    logger.info(
        "document_extract filename=%s method=%s chars=%s",
        filename,
        method,
        len(text),
    )
    return DocumentExtractResult(
        text=text,
        filename=filename,
        method=method,
        char_count=len(text),
    )
